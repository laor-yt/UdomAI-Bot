import os
import fitz # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
import speech_recognition as sr
from pydub import AudioSegment
import imageio_ffmpeg
import uuid

# Configure pydub and system environment to use the ffmpeg binary from imageio_ffmpeg
ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
ffmpeg_dir = os.path.dirname(ffmpeg_exe)
if ffmpeg_dir not in os.environ.get("PATH", ""):
    os.environ["PATH"] = ffmpeg_dir + os.pathsep + os.environ.get("PATH", "")

AudioSegment.converter = ffmpeg_exe
AudioSegment.ffmpeg = ffmpeg_exe
AudioSegment.ffprobe = os.path.join(ffmpeg_dir, "ffprobe.exe") if os.name == 'nt' else os.path.join(ffmpeg_dir, "ffprobe")

# Maximum number of characters to extract to prevent AI context overflow
MAX_CHARS = 12000

def parse_document(file_path: str, mime_type: str) -> str:
    """Extracts text from a document file based on its MIME type."""
    ext = os.path.splitext(file_path)[1].lower()
    
    text = ""
    try:
        if ext == ".pdf" or "pdf" in mime_type:
            doc = fitz.open(file_path)
            for page in doc:
                text += page.get_text() + "\n"
                if len(text) > MAX_CHARS:
                    break
                    
        elif ext in [".docx", ".doc"] or "word" in mime_type:
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            
        elif ext in [".xlsx", ".xls", ".csv"] or "excel" in mime_type or "spreadsheet" in mime_type:
            if ext == ".csv":
                df = pd.read_csv(file_path)
            else:
                df = pd.read_excel(file_path)
            text = df.head(100).to_string() # Read first 100 rows
            
        elif ext in [".pptx", ".ppt"] or "powerpoint" in mime_type or "presentation" in mime_type:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) > MAX_CHARS:
                    break
        else:
            return "Unsupported document format for text extraction."
            
    except Exception as e:
        print(f"Error parsing document: {e}")
        return f"Error reading file: {e}"
        
    if not text.strip():
        return "The document appears to be empty or contains no readable text (it might be scanned images)."
        
    return text[:MAX_CHARS]


def transcribe_audio_video(file_path: str, src_lang: str = "auto") -> str:
    """Extracts audio and transcribes 100% of spoken speech for media in ANY global language (Chinese, Khmer, English, Japanese, Korean, Vietnamese, Thai, French, Spanish, German, Russian, etc.)."""
    temp_wav = f"temp_{uuid.uuid4().hex}.wav"
    temp_mp3 = f"temp_{uuid.uuid4().hex}.mp3"
    
    LANG_CODE_MAP = {
        "zh": "zh-CN", "km": "km-KH", "en": "en-US", "vi": "vi-VN", "th": "th-TH",
        "ja": "ja-JP", "ko": "ko-KR", "fr": "fr-FR", "es": "es-ES", "de": "de-DE",
        "ru": "ru-RU", "ar": "ar-SA", "hi": "hi-IN", "id": "id-ID", "pt": "pt-PT", "it": "it-IT"
    }
    
    try:
        # Extract audio track to lightweight 24kHz mono MP3 format using imageio_ffmpeg (immune to missing ffprobe)
        ffmpeg_bin = imageio_ffmpeg.get_ffmpeg_exe()
        import ffmpeg
        (
            ffmpeg
            .input(file_path)
            .output(temp_mp3, acodec='libmp3lame', ar='24000', ac=1, ab='64k')
            .overwrite_output()
            .run(cmd=ffmpeg_bin, capture_stdout=True, capture_stderr=True)
        )
        
        if not os.path.exists(temp_mp3) or os.path.getsize(temp_mp3) == 0:
            return "ERROR: Failed to extract audio track from media file."
        
        SRC_LANG_NAMES = {
            "zh": "Mandarin Chinese (中文 / 简体中文 / 繁體中文)",
            "km": "Khmer (ភាសាខ្មែរ)",
            "en": "English",
            "vi": "Vietnamese (Tiếng Việt)",
            "th": "Thai (ภาษาไทย)",
            "ja": "Japanese (日本語)",
            "ko": "Korean (한국어)",
            "fr": "French (Français)",
            "es": "Spanish (Español)",
            "de": "German (Deutsch)",
            "ru": "Russian (Русский)",
            "ar": "Arabic (العربية)",
            "hi": "Hindi (हिन्दी)",
            "id": "Indonesian (Bahasa Indonesia)"
        }
        
        clean_src = (src_lang or "auto").lower().strip()[:2]
        if clean_src == "zh":
            transcription_instruction = (
                "STRICT VERBATIM CHINESE MANDATE: The audio is in Chinese (Mandarin / Cantonese / 汉语 / 普通话 / 粤语).\n"
                "Transcribe 100% of ALL spoken Chinese words ONLY in standard Chinese Hanzi characters (简体字 or 繁體字) with 100% verbatim precision.\n"
                "RULES:\n"
                "1. Output exact Chinese spoken words in Hanzi characters. Do NOT translate into English, Khmer, or any other language.\n"
                "2. Do NOT omit or change any Chinese word. Keep exact original sentence order.\n"
                "3. If a word or phrase is completely inaudible or noisy, write '[inaudible]' instead of inventing words.\n"
                "4. Output ONLY the raw verbatim Chinese text script with NO commentary, NO pinyin, and NO translation."
            )
        elif clean_src in SRC_LANG_NAMES:
            target_lang_desc = SRC_LANG_NAMES[clean_src]
            transcription_instruction = (
                f"STRICT VERBATIM MANDATE: The audio is in {target_lang_desc}.\n"
                f"Transcribe 100% of ALL spoken words ONLY in {target_lang_desc} ({clean_src}) with 100% verbatim precision.\n"
                f"RULES:\n"
                f"1. If speaker says 'apple', write 'apple' — NOT 'banana'. Do NOT translate into English or substitute any word.\n"
                f"2. If a word or phrase is completely inaudible, noisy, or muffled, write '[inaudible]' instead of guessing wrong words.\n"
                f"3. Output ONLY raw verbatim transcript in {target_lang_desc} script with NO commentary."
            )
        else:
            transcription_instruction = (
                "STRICT VERBATIM MANDATE: Transcribe 100% of ALL spoken words in this audio in its ORIGINAL spoken language "
                "(Chinese 中文, Khmer ភាសាខ្មែរ, English, Vietnamese, Thai, Korean, Japanese, French, Spanish, German, Russian, etc.) with exact precision.\n"
                "RULES:\n"
                "1. If speaker says 'apple', write 'apple' — NOT 'banana'. Do NOT translate or substitute any word.\n"
                "2. If a word or phrase is completely inaudible, noisy, or muffled, write '[inaudible]' instead of guessing wrong words.\n"
                "3. Output ONLY raw verbatim transcript in original spoken language script with NO commentary."
            )

        # 1. Try Gemini 2.0 Flash first, then 1.5 Flash fallback
        api_key = os.environ.get("GEMINI_API_KEY")
        if api_key and os.path.exists(temp_mp3):
            models_to_try = ["gemini-2.0-flash", "gemini-1.5-flash", "gemini-1.5-pro"]
            for gem_model in models_to_try:
                try:
                    import base64
                    import requests
                    with open(temp_mp3, "rb") as f:
                        b64 = base64.b64encode(f.read()).decode("utf-8")
                    url = f"https://generativelanguage.googleapis.com/v1beta/models/{gem_model}:generateContent?key={api_key}"
                    payload = {
                        "contents": [{
                            "parts": [
                                {"text": transcription_instruction},
                                {"inline_data": {"mime_type": "audio/mp3", "data": b64}}
                            ]
                        }],
                        "generationConfig": {"temperature": 0.0, "topP": 0.05, "topK": 1}
                    }
                    res = requests.post(url, json=payload, timeout=90).json()
                    if "candidates" in res and res["candidates"]:
                        p_text = res["candidates"][0]["content"]["parts"][0]["text"].strip()
                        if p_text and len(p_text) > 3:
                            return p_text
                except Exception as e:
                    print(f"Gemini Model ({gem_model}) Error: {e}")

        # 2. Multi-language Chunked Google SpeechRecognition Fallback
        (
            ffmpeg
            .input(temp_mp3)
            .output(temp_wav, ar='16000', ac=1)
            .overwrite_output()
            .run(cmd=ffmpeg_bin, capture_stdout=True, capture_stderr=True)
        )
        recognizer = sr.Recognizer()
        full_transcript = []
        
        if clean_src == "zh":
            fallback_langs = ["zh-CN", "zh-TW", "zh-HK"]
        elif clean_src in LANG_CODE_MAP:
            primary_lang = LANG_CODE_MAP[clean_src]
            fallback_langs = [primary_lang] + [l for l in ["zh-CN", "en-US", "km-KH", "ja-JP", "vi-VN"] if l != primary_lang]
        else:
            fallback_langs = ["en-US", "km-KH", "zh-CN", "vi-VN", "th-TH", "ja-JP", "ko-KR"]

        with sr.AudioFile(temp_wav) as source:
            audio_data = recognizer.record(source)
            for l_try in fallback_langs:
                try:
                    chunk_text = recognizer.recognize_google(audio_data, language=l_try)
                    if chunk_text:
                        full_transcript.append(chunk_text)
                        break
                except Exception:
                    continue
                    
        result_text = " ".join(full_transcript).strip()
        if result_text:
            return result_text
            
        return "Could not recognize speech in media."
    except Exception as e:
        print(f"Error transcribing media: {e}")
        return f"Error extracting speech: {e}"
    finally:
        if os.path.exists(temp_wav): os.remove(temp_wav)
        if os.path.exists(temp_mp3): os.remove(temp_mp3)
